#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Standalone OCR + file ingestion pipeline.
- Walks input directories
- Extracts text from supported formats
- OCRs images (and images embedded in XLSX; PDFs fallback to OCR if needed)
- Emits per-file .txt and .meta.json outputs under --out-dir preserving relative paths
- Writes manifest.jsonl and manifest.csv for review
- Skips unchanged files using sha256 + mtime
"""
from providers_db import SessionLocal, Provider, ProviderAlias, ProviderNote
from sqlalchemy import select
import argparse, os, sys, io, json, csv, shutil, hashlib, time
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional

# 3rd-party (installed via pip)
import chardet

# Lazy imports inside functions for heavy libs (PyPDF2, pdf2image, openpyxl, PIL, pytesseract, docx, pptx)

SUPPORTED_EXTS = {".txt", ".md", ".json", ".csv",
                  ".pdf", ".docx", ".pptx",
                  ".xlsx",
                  ".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp"}

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp"}

def sha256_of_file(p: Path, blocksize: int = 1024 * 1024) -> str:
    h = hashlib.sha256()
    with p.open("rb") as f:
        for chunk in iter(lambda: f.read(blocksize), b""):
            h.update(chunk)
    return h.hexdigest()

def safe_text(text: str) -> str:
    # Normalize line endings; keep it simple
    return text.replace("\r\n", "\n").replace("\r", "\n")

def read_text_guess_encoding(p: Path) -> str:
    data = p.read_bytes()
    guess = chardet.detect(data)
    enc = guess.get("encoding") or "utf-8"
    try:
        return data.decode(enc, errors="ignore")
    except Exception:
        return data.decode("utf-8", errors="ignore")

def record_alias(abbr, alias, confidence, source):
    db = SessionLocal()
    try:
        p = db.execute(select(Provider).where(Provider.abbr == abbr)).scalar_one_or_none()
        if not p: return
        norm = alias.strip().lower()
        exists = db.execute(select(ProviderAlias).where(
            ProviderAlias.provider_id == p.id,
            ProviderAlias.normalized == norm
        )).scalar_one_or_none()
        if not exists:
            db.add(ProviderAlias(provider_id=p.id, alias=alias, normalized=norm,
                                 confidence=confidence, source=source, pending_review=True))
            db.commit()
    finally:
        db.close()

def record_note(abbr, note, source):
    db = SessionLocal()
    try:
        p = db.execute(select(Provider).where(Provider.abbr == abbr)).scalar_one_or_none()
        if not p: return
        db.add(ProviderNote(provider_id=p.id, note=note, source=source))
        db.commit()
    finally:
        db.close()

def extract_txt_md_json_csv(p: Path) -> str:
    suf = p.suffix.lower()
    if suf in {".txt", ".md"}:
        return read_text_guess_encoding(p)
    if suf == ".json":
        try:
            obj = json.loads(read_text_guess_encoding(p))
            return json.dumps(obj, indent=2, ensure_ascii=False)
        except Exception:
            return read_text_guess_encoding(p)
    if suf == ".csv":
        # Keep original CSV but trim absurdly long lines/rows to be safe
        return read_text_guess_encoding(p)
    return ""

def extract_docx(p: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(p))
        return "\n".join(par.text for par in doc.paragraphs)
    except Exception as e:
        return f"(docx extract failed: {e})"

def extract_pptx(p: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
                except Exception:
                    continue
        return "\n".join(parts)
    except Exception as e:
        return f"(pptx extract failed: {e})"

def extract_pdf(p: Path, tesseract_cmd: Optional[str]) -> (str, Dict[str, Any]):
    meta = {"pages": 0, "ocr_used": False}
    text = ""
    # Try native text
    try:
        import PyPDF2
        with p.open("rb") as f:
            reader = PyPDF2.PdfReader(f)
            meta["pages"] = len(reader.pages)
            pages = []
            for page in reader.pages:
                try:
                    pages.append(page.extract_text() or "")
                except Exception:
                    pages.append("")
            text = "\n".join(pages)
    except Exception:
        text = ""

    if text.strip():
        return text, meta

    # Fallback to OCR with pdf2image
    try:
        from pdf2image import convert_from_path
        from PIL import Image
        import pytesseract
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        images = convert_from_path(str(p))
        meta["pages"] = len(images)
        ocr_pages = []
        for im in images:
            ocr_pages.append(pytesseract.image_to_string(im))
        meta["ocr_used"] = True
        return "\n".join(ocr_pages), meta
    except Exception as e:
        return f"(pdf ocr failed: {e})", meta

def ocr_image_bytes(data: bytes, tesseract_cmd: Optional[str]) -> str:
    try:
        from PIL import Image
        import pytesseract
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        im = Image.open(io.BytesIO(data))
        return pytesseract.image_to_string(im)
    except Exception as e:
        return f"(image ocr failed: {e})"

def ocr_image_file(p: Path, tesseract_cmd: Optional[str]) -> str:
    try:
        data = p.read_bytes()
        return ocr_image_bytes(data, tesseract_cmd)
    except Exception as e:
        return f"(image read failed: {e})"

def extract_xlsx_all(p: Path, tesseract_cmd: Optional[str]) -> (str, Dict[str, Any]):
    """
    - All sheets → TSV-like text
    - All images in xl/media → OCR
    """
    meta = {"sheets": 0, "images_ocr": 0, "ocr_used": False}
    parts = []

    # Sheets
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filename=str(p), data_only=True, read_only=True)
        meta["sheets"] = len(wb.worksheets)
        for ws in wb.worksheets:
            parts.append(f"## Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                line = "\t".join("" if v is None else str(v) for v in row)
                if line.strip():
                    parts.append(line)
    except Exception as e:
        parts.append(f"(sheet read failed: {e})")

    # Images under xl/media/*
    try:
        import zipfile
        with zipfile.ZipFile(str(p), "r") as zf:
            media_names = [n for n in zf.namelist() if n.lower().startswith("xl/media/")]
            if media_names:
                parts.append("\n### OCR (images in workbook)\n")
            for name in media_names:
                try:
                    data = zf.read(name)
                    txt = ocr_image_bytes(data, tesseract_cmd)
                    txt = (txt or "").strip()
                    if txt:
                        parts.append(f"[image {Path(name).name}]\n{txt}\n")
                        meta["images_ocr"] += 1
                        meta["ocr_used"] = True
                except Exception:
                    continue
    except Exception:
        pass

    return "\n".join(parts).strip(), meta

def extract_any(p: Path, tesseract_cmd: Optional[str]) -> (str, Dict[str, Any]):
    suf = p.suffix.lower()
    meta: Dict[str, Any] = {"ocr_used": False}
    if suf in {".txt", ".md", ".json", ".csv"}:
        return extract_txt_md_json_csv(p), meta
    if suf == ".docx":
        return extract_docx(p), meta
    if suf == ".pptx":
        return extract_pptx(p), meta
    if suf == ".pdf":
        text, m = extract_pdf(p, tesseract_cmd)
        meta.update(m)
        return text, meta
    if suf == ".xlsx":
        text, m = extract_xlsx_all(p, tesseract_cmd)
        meta.update(m)
        return text, meta
    if suf in IMAGE_EXTS:
        txt = ocr_image_file(p, tesseract_cmd)
        meta["ocr_used"] = True
        return txt, meta
    return "", meta

def write_text_and_meta(out_root: Path, src_root: Path, src_path: Path,
                        text: str, meta: Dict[str, Any],
                        sha256: str, mtime: float, size: int) -> Dict[str, Any]:
    rel = src_path.relative_to(src_root)
    stem = rel.as_posix()
    out_txt = out_root / rel.with_suffix(rel.suffix + ".txt")
    out_meta = out_root / rel.with_suffix(rel.suffix + ".meta.json")

    out_txt.parent.mkdir(parents=True, exist_ok=True)

    out_txt.write_text(safe_text(text), encoding="utf-8")
    meta_full = {
        "src_rel_path": stem,
        "src_abs_path": str(src_path),
        "size": size,
        "mtime": mtime,
        "mtime_iso": datetime.fromtimestamp(mtime).isoformat(),
        "sha256": sha256,
        "ext": src_path.suffix.lower(),
        **meta
    }
    out_meta.write_text(json.dumps(meta_full, indent=2, ensure_ascii=False), encoding="utf-8")
    return {"rel": stem, **meta_full, "out_txt": str(out_txt), "out_meta": str(out_meta)}

def should_skip(out_root: Path, src_root: Path, src_path: Path, sha256: str, mtime: float) -> bool:
    rel = src_path.relative_to(src_root)
    out_meta = out_root / rel.with_suffix(rel.suffix + ".meta.json")
    if not out_meta.exists():
        return False
    try:
        prev = json.loads(out_meta.read_text(encoding="utf-8"))
        return prev.get("sha256") == sha256 and abs(prev.get("mtime", 0) - mtime) < 0.5
    except Exception:
        return False

def walk_inputs(input_dirs: List[Path]) -> List[Path]:
    files = []
    for root in input_dirs:
        for p in root.rglob("*"):
            if p.is_file() and p.suffix.lower() in SUPPORTED_EXTS and not p.name.startswith("~$"):
                files.append(p)
    return files

def main():
    ap = argparse.ArgumentParser(description="OCR + File ingestion (review-first).")
    ap.add_argument("--in-dir", action="append", required=True,
                    help="Input directory (can specify multiple)")
    ap.add_argument("--out-dir", required=True, help="Output directory")
    ap.add_argument("--tesseract-cmd", default=None,
                    help="Path to tesseract binary (e.g., /opt/homebrew/bin/tesseract)")
    ap.add_argument("--force", action="store_true", help="Reprocess even if unchanged")
    ap.add_argument("--manifest-csv", default="manifest.csv", help="CSV manifest filename")
    ap.add_argument("--manifest-jsonl", default="manifest.jsonl", help="JSONL manifest filename")
    args = ap.parse_args()

    out_root = Path(args.out_dir).expanduser().resolve()
    out_root.mkdir(parents=True, exist_ok=True)
    inputs = [Path(d).expanduser().resolve() for d in args.in_dir]
    for d in inputs:
        if not d.exists():
            sys.stderr.write(f"[warn] Input dir missing: {d}\n")

    files = walk_inputs([d for d in inputs if d.exists()])
    if not files:
        print("[info] No files found.")
        return

    manifest_rows: List[Dict[str, Any]] = []
    jsonl_path = out_root / args.manifest_jsonl
    csv_path = out_root / args.manifest_csv

    # (Re)create manifests
    if jsonl_path.exists():
        jsonl_path.unlink()
    if csv_path.exists():
        csv_path.unlink()

    t0 = time.time()
    processed = 0

    for f in files:
        try:
            src_root = next((d for d in inputs if str(f).startswith(str(d))), inputs[0])
            size = f.stat().st_size
            mtime = f.stat().st_mtime
            h = sha256_of_file(f)

            if not args.force and should_skip(out_root, src_root, f, h, mtime):
                # Load meta quickly to include in manifest
                rel = f.relative_to(src_root).as_posix()
                meta_path = out_root / f.relative_to(src_root).with_suffix(f.suffix + ".meta.json")
                try:
                    prev = json.loads(meta_path.read_text(encoding="utf-8"))
                    prev["rel"] = rel
                    prev["out_txt"] = str(out_root / f.relative_to(src_root).with_suffix(f.suffix + ".txt"))
                    prev["out_meta"] = str(meta_path)
                    manifest_rows.append(prev)
                except Exception:
                    pass
                continue

            text, meta = extract_any(f, args.tesseract_cmd)
            rec = write_text_and_meta(out_root, src_root, f, text, meta, h, mtime, size)
            manifest_rows.append(rec)
            processed += 1
            print(f"[ok] {f}  (ocr={rec.get('ocr_used')})")
        except Exception as e:
            # Record a failed record with error
            try:
                src_root = next((d for d in inputs if str(f).startswith(str(d))), inputs[0])
            except Exception:
                src_root = inputs[0]
            rel = f.relative_to(src_root).as_posix()
            manifest_rows.append({
                "rel": rel,
                "src_rel_path": rel,
                "src_abs_path": str(f),
                "size": f.stat().st_size if f.exists() else None,
                "mtime": f.stat().st_mtime if f.exists() else None,
                "mtime_iso": datetime.fromtimestamp(f.stat().st_mtime).isoformat() if f.exists() else None,
                "sha256": None,
                "ext": f.suffix.lower(),
                "ocr_used": None,
                "error": str(e)
            })
            print(f"[err] {f}: {e}", file=sys.stderr)

    # Write JSONL + CSV
    with (out_root / args.manifest_jsonl).open("w", encoding="utf-8") as jf:
        for row in manifest_rows:
            jf.write(json.dumps(row, ensure_ascii=False) + "\n")

    if manifest_rows:
        fieldnames = sorted(set().union(*[row.keys() for row in manifest_rows]))
        with (out_root / args.manifest_csv).open("w", newline="", encoding="utf-8") as cf:
            w = csv.DictWriter(cf, fieldnames=fieldnames)
            w.writeheader()
            for row in manifest_rows:
                w.writerow(row)

    dt = time.time() - t0
    print(f"[done] processed={processed}/{len(files)} in {dt:.1f}s")
    print(f"[out] {jsonl_path}")
    print(f"[out] {csv_path}")

if __name__ == "__main__":
    main()